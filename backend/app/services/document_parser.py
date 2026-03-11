"""Document parsing service for PDF and PPTX files."""

import io
import logging
from dataclasses import dataclass

import httpx

logger = logging.getLogger(__name__)

MAX_DOWNLOAD_SIZE = 20 * 1024 * 1024  # 20MB


@dataclass
class ParseResult:
    """Result of document parsing."""

    text: str
    doc_type: str
    pages: int
    success: bool
    error: str | None = None


class DocumentParserService:
    """Download and parse PDF/PPTX documents."""

    @staticmethod
    def is_document_url(url: str) -> str | None:
        """Detect if URL points to a document. Returns 'pdf', 'pptx', or None."""
        lower = url.lower().split("?")[0]
        if lower.endswith(".pdf"):
            return "pdf"
        if lower.endswith(".pptx"):
            return "pptx"
        return None

    async def download_and_parse(self, url: str, timeout: float = 30.0) -> ParseResult:
        """Download a document and extract text.

        Args:
            url: URL of the document.
            timeout: Download timeout in seconds.

        Returns:
            ParseResult with extracted text or error.
        """
        doc_type = self.is_document_url(url)
        if not doc_type:
            return ParseResult(text="", doc_type="", pages=0, success=False, error="Not a supported document URL")

        try:
            data = await self._download(url, timeout)
        except Exception as e:
            return ParseResult(text="", doc_type=doc_type, pages=0, success=False, error=str(e))

        if doc_type == "pdf":
            return self._parse_pdf(data)
        elif doc_type == "pptx":
            return self._parse_pptx(data)
        return ParseResult(text="", doc_type=doc_type, pages=0, success=False, error=f"Unsupported type: {doc_type}")

    async def _download(self, url: str, timeout: float) -> bytes:
        """Download document with size limit check."""
        async with httpx.AsyncClient(
            timeout=timeout,
            follow_redirects=True,
            headers={"User-Agent": "Mozilla/5.0 (compatible; AINewsRadio/1.0)"},
        ) as client:
            response = await client.get(url)
            response.raise_for_status()

        if len(response.content) > MAX_DOWNLOAD_SIZE:
            raise ValueError(f"File too large: {len(response.content)} bytes (max {MAX_DOWNLOAD_SIZE})")

        return response.content

    def _parse_pdf(self, data: bytes) -> ParseResult:
        """Parse PDF document."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            pages_text = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            full_text = "\n\n".join(pages_text)
            return ParseResult(text=full_text, doc_type="pdf", pages=len(reader.pages), success=True)
        except Exception as e:
            return ParseResult(text="", doc_type="pdf", pages=0, success=False, error=str(e))

    def _parse_pptx(self, data: bytes) -> ParseResult:
        """Parse PPTX document."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            slides_text = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                if texts:
                    slides_text.append("\n".join(texts))
            full_text = "\n\n".join(slides_text)
            return ParseResult(text=full_text, doc_type="pptx", pages=len(prs.slides), success=True)
        except Exception as e:
            return ParseResult(text="", doc_type="pptx", pages=0, success=False, error=str(e))
